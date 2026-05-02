"""High-level builder for coursework slide decks.

Adapted from references/mckinsey-pptx/mckinsey_pptx/builder.py (MIT License).
Simplified for Chinese coursework: fewer templates, slide card integration.
"""
from __future__ import annotations
from typing import Sequence, Optional, Dict, Any, Iterable

from .base import init_presentation
from .theme import Theme, DEFAULT_THEME
from .slides.cover import add_cover
from .slides.content import add_bullet_list, add_two_column, add_question
from .slides.section import add_section_divider, add_stat_hero
from .slides.comparison import add_pros_cons, add_before_after
from .slides.timeline import add_process_flow, add_phases
from .slides.summary import add_three_takeaways, add_closing, add_quote
from .slides.visual import add_diagram, add_icon_grid, add_table_slide


# --- Layout -> template mapping ---
# slide-writer.md uses these layout names; map them to engine templates.
_LAYOUT_MAP = {
    "title-hero": "cover",
    "title-with-subtitle": "cover",
    "question": "question",
    "contrast": "question",
    "problem-card": "question",
    "diagram": "diagram",
    "two-column": "two_column",
    "process-flow": "process_flow",
    "before-after": "before_after",
    "example-card": "bullet_list",
    "quote-plus-analysis": "quote",
    "section-divider": "section_divider",
    "big-number": "stat_hero",
    "three-takeaways": "three_takeaways",
    "closing-statement": "closing",
    "icon-grid": "icon_grid",
}

def _process_flow_adapter(prs, *, steps=None, **kwargs):
    """Adapter for process_flow that accepts 'steps' as list of strings."""
    return add_process_flow(prs, steps=steps or [], **kwargs)


# --- Template registry ---
_REGISTRY = {
    "cover": add_cover,
    "bullet_list": add_bullet_list,
    "two_column": add_two_column,
    "question": add_question,
    "section_divider": add_section_divider,
    "stat_hero": add_stat_hero,
    "pros_cons": add_pros_cons,
    "before_after": add_before_after,
    "process_flow": _process_flow_adapter,
    "phases": add_phases,
    "three_takeaways": add_three_takeaways,
    "closing": add_closing,
    "quote": add_quote,
    "diagram": add_diagram,
    "icon_grid": add_icon_grid,
    "table": add_table_slide,
}


def infer_slide_type(spec: Dict[str, Any]) -> str:
    """Pick a template by examining the spec's keys."""
    # Structural
    if "course" in spec and "student" in spec:
        return "cover"
    if "section_number" in spec and "section_title" in spec:
        return "section_divider"
    if "stat" in spec and "stat_label" in spec:
        return "stat_hero"
    if "quote" in spec and "author" in spec:
        return "quote"

    # Content patterns
    if "pros" in spec and "cons" in spec:
        return "pros_cons"
    if "left_items" in spec and "right_items" in spec:
        return "before_after"
    if "steps" in spec:
        return "process_flow"
    if "phases" in spec:
        return "phases"
    if "takeaways" in spec:
        return "three_takeaways"
    if "headers" in spec and "rows" in spec:
        return "table"
    if "items" in spec and isinstance(spec["items"], list):
        # icon_grid items are dicts with "title" key; diagram items are strings
        if spec["items"] and isinstance(spec["items"][0], dict):
            return "icon_grid"
        return "diagram"
    if "bullets" in spec:
        return "bullet_list"
    if "question" in spec:
        return "question"
    if "message" in spec and len(spec) <= 3:
        return "closing"

    # Default
    return "bullet_list"


class PresentationBuilder:
    """Compose a deck of coursework-style slides."""

    def __init__(self, theme: Theme = DEFAULT_THEME, *,
                 auto_page_numbers: bool = True):
        self.theme = theme
        self.auto_page_numbers = auto_page_numbers
        self.prs = init_presentation(theme)
        self._page = 0

    def add(self, slide_type: str, **kwargs):
        fn = _REGISTRY.get(slide_type)
        if fn is None:
            raise ValueError(
                f"Unknown slide type: {slide_type}. "
                f"Available: {sorted(_REGISTRY)}")
        if self.auto_page_numbers:
            self._page += 1
            kwargs.setdefault("page_number", self._page)
        kwargs.setdefault("theme", self.theme)
        return fn(self.prs, **kwargs)

    def add_spec(self, spec: Dict[str, Any]):
        spec = dict(spec)
        slide_type = spec.pop("type", None) or infer_slide_type(spec)
        speaker_notes = spec.pop("_speaker_notes", None)

        # Map layout alias if present
        layout = spec.pop("layout", None)
        if layout and slide_type == "bullet_list":
            mapped = _LAYOUT_MAP.get(layout)
            if mapped and mapped != "bullet_list":
                slide_type = mapped

        slide = self.add(slide_type, **spec)

        # Apply speaker notes after slide creation
        if speaker_notes:
            from pptx.oxml.ns import qn
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = speaker_notes

        return slide

    def add_specs(self, specs: Iterable[Dict[str, Any]]):
        return [self.add_spec(s) for s in specs]

    def save(self, path: str):
        self.prs.save(path)
        return path


def build_from_slides_md(slides_md_path: str, output_path: str, *,
                         theme: Theme = DEFAULT_THEME) -> str:
    """Parse a slide card markdown file and build a PPTX."""
    from .slide_card_parser import parse_slides_md
    specs = parse_slides_md(slides_md_path)
    b = PresentationBuilder(theme=theme)
    b.add_specs(specs)
    return b.save(output_path)
