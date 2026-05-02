"""Common slide primitives: title, underline, footer, text/shape helpers.

Adapted from references/mckinsey-pptx/mckinsey_pptx/base.py (MIT License).
"""
from __future__ import annotations
from typing import Optional

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from .theme import Theme, DEFAULT_THEME


def enable_text_shrink(text_frame):
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE_FIT
    except Exception:
        pass


# ---------- text helpers ----------

def set_run(run, text, *, size=None, bold=False, italic=False, color=None,
            family=None):
    run.text = text
    if size is not None:
        run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if family:
        run.font.name = family
    if color is not None:
        run.font.color.rgb = color


def add_textbox(slide, left_in, top_in, width_in, height_in, *,
                fill=None, line=None, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                  Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if line is None:
        tb.line.fill.background()
    return tb


def write_paragraph(tf, text, *, size, bold=False, italic=False, color=None,
                    family="Microsoft YaHei", align=PP_ALIGN.LEFT,
                    space_before=None, space_after=None, level=0,
                    bullet=False, first=False):
    if first and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_before is not None:
        p.space_before = Pt(space_before)
    if space_after is not None:
        p.space_after = Pt(space_after)
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, italic=italic, color=color,
            family=family)
    if bullet:
        _set_bullet(p, color)
    else:
        _clear_bullet(p)
    return p


def _set_bullet(paragraph, color: Optional[RGBColor]):
    pPr = (paragraph._pPr if paragraph._pPr is not None
           else paragraph._p.get_or_add_pPr())
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    bu = etree.SubElement(pPr, qn("a:buChar"))
    bu.set("char", "\u2022")
    pPr.set("indent", "-228600")
    pPr.set("marL", "228600")


def _clear_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    none = pPr.find(qn("a:buNone"))
    if none is None:
        etree.SubElement(pPr, qn("a:buNone"))


# ---------- shape helpers ----------

def add_rect(slide, left_in, top_in, width_in, height_in, *,
             fill=None, line=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(left_in), Inches(top_in),
                               Inches(width_in), Inches(height_in))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_width is not None:
            s.line.width = Pt(line_width)
    s.text_frame.margin_left = Emu(0)
    s.text_frame.margin_right = Emu(0)
    s.text_frame.margin_top = Emu(0)
    s.text_frame.margin_bottom = Emu(0)
    return s


def add_oval(slide, left_in, top_in, width_in, height_in, *,
             fill=None, line=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(left_in), Inches(top_in),
                               Inches(width_in), Inches(height_in))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_width is not None:
            s.line.width = Pt(line_width)
    return s


def add_line(slide, x1_in, y1_in, x2_in, y2_in, *,
             color, width_pt=0.75, dash=None):
    line = slide.shapes.add_connector(1, Inches(x1_in), Inches(y1_in),
                                      Inches(x2_in), Inches(y2_in))
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    if dash is not None:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        line.line.dash_style = dash
    return line


# ---------- decoration primitives ----------

def add_icon_circle(slide, cx_in, cy_in, diameter_in, text, *,
                    fill=None, text_color=None, font_size=14,
                    theme: Theme = DEFAULT_THEME):
    """A filled circle with centered text (number, symbol, or emoji).

    Returns the oval shape. Text is centered both horizontally and vertically.
    """
    pal = theme.palette
    r = diameter_in / 2
    oval = add_oval(slide, cx_in - r, cy_in - r, diameter_in, diameter_in,
                    fill=fill or pal.primary, line=None)
    tf = oval.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, text, size=font_size, bold=True,
            color=text_color or pal.white, family=theme.typography.family)
    return oval


def add_accent_bar(slide, left_in, top_in, width_in, height_in, *,
                   color=None, theme: Theme = DEFAULT_THEME):
    """A thin decorative color bar (stripe, divider, or card top accent).

    Returns the rectangle shape.
    """
    return add_rect(slide, left_in, top_in, width_in, height_in,
                    fill=color or theme.palette.primary)


def add_label_badge(slide, left_in, top_in, text, *,
                    fill=None, text_color=None, font_size=10,
                    padding_x=0.15, theme: Theme = DEFAULT_THEME):
    """A rounded-rectangle badge with text inside.

    Width is auto-calculated from text length. Returns the shape.
    """
    pal = theme.palette
    # Approximate width: 0.12" per character + padding
    char_w = 0.13 if any(ord(c) > 127 for c in text) else 0.075
    w = len(text) * char_w + padding_x * 2
    h = 0.3
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(left_in), Inches(top_in),
                               Inches(w), Inches(h))
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = fill or pal.primary
    s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, text, size=font_size, bold=True,
            color=text_color or pal.white, family=theme.typography.family)
    return s


# ---------- slide chrome ----------

def add_title(slide, text, theme: Theme = DEFAULT_THEME, *,
              with_underline=True):
    layout = theme.layout
    pal = theme.palette
    typo = theme.typography
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    tb = add_textbox(slide, layout.margin_left_in, layout.title_top_in,
                     width, layout.title_height_in)
    write_paragraph(tb.text_frame, text, size=typo.title_size, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)
    enable_text_shrink(tb.text_frame)
    if with_underline:
        add_line(slide, layout.margin_left_in, layout.title_underline_top_in,
                 layout.slide_width_in - layout.margin_right_in,
                 layout.title_underline_top_in,
                 color=pal.rule_gray, width_pt=0.75)
    return tb


def add_footer(slide, theme: Theme = DEFAULT_THEME, *,
               page_number=None, source=None, footnote=None,
               copyright_text=None):
    layout = theme.layout
    pal = theme.palette
    typo = theme.typography

    foot_y = layout.footer_top_in
    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in

    add_line(slide, layout.margin_left_in, foot_y,
             layout.slide_width_in - layout.margin_right_in, foot_y,
             color=pal.rule_gray, width_pt=0.5)

    y = foot_y + 0.05
    if footnote:
        tb = add_textbox(slide, layout.margin_left_in, y, width / 2, 0.18)
        write_paragraph(tb.text_frame, footnote, size=typo.footer_size,
                        color=pal.text_dark, family=typo.family, first=True)
        y += 0.18
    if source:
        tb = add_textbox(slide, layout.margin_left_in, y, width / 2, 0.18)
        write_paragraph(tb.text_frame, f"Source: {source}",
                        size=typo.footer_size, color=pal.text_dark,
                        family=typo.family, first=True)

    right_w = 4.0
    right_left = layout.slide_width_in - layout.margin_right_in - right_w
    cp = copyright_text if copyright_text is not None else theme.copyright_text
    parts = []
    if cp:
        parts.append(cp)
    if page_number is not None:
        if parts:
            parts.append("   ")
        parts.append(str(page_number))
    if parts:
        tb = add_textbox(slide, right_left, foot_y + 0.1, right_w, 0.2)
        write_paragraph(tb.text_frame, "".join(parts),
                        size=typo.footer_size, color=pal.footer_gray,
                        family=typo.family, align=PP_ALIGN.RIGHT, first=True)


def init_presentation(theme: Theme = DEFAULT_THEME):
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width = Inches(theme.layout.slide_width_in)
    prs.slide_height = Inches(theme.layout.slide_height_in)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_chrome(slide, *, title, theme=DEFAULT_THEME, page_number=None,
               source=None, footnote=None, section_marker=None,
               with_underline=True):
    add_title(slide, title, theme, with_underline=with_underline)
    add_footer(slide, theme, page_number=page_number, source=source,
               footnote=footnote)
