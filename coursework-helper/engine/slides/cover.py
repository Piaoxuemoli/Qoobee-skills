"""Cover slide — title page for coursework presentations."""
from __future__ import annotations
from typing import Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..base import (
    blank_slide, add_rect, add_line, add_textbox, write_paragraph,
    enable_text_shrink, add_title, add_footer, add_accent_bar,
)
from ..theme import Theme, DEFAULT_THEME


def add_cover(prs, *,
              title: str,
              subtitle: Optional[str] = None,
              course: Optional[str] = None,
              student: Optional[str] = None,
              date: Optional[str] = None,
              page_number=None,
              theme: Theme = DEFAULT_THEME):
    """First-page cover: large title, subtitle, course/student/date metadata.

    Includes accent stripe, decorative lines, and bottom metadata block
    to reduce the gap between subtitle and metadata.
    """
    slide = blank_slide(prs)

    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Accent stripe on right edge
    stripe_w = 0.35
    add_rect(slide, layout.slide_width_in - stripe_w, 0,
             stripe_w, layout.slide_height_in, fill=pal.primary)

    # Big title
    title_left = layout.margin_left_in + 0.15
    title_top = 2.0
    title_w = layout.slide_width_in - layout.margin_left_in - 1.2
    tb = add_textbox(slide, title_left, title_top, title_w, 1.6)
    write_paragraph(tb.text_frame, title,
                    size=typo.title_size + 18, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)
    enable_text_shrink(tb.text_frame)

    # Subtitle
    if subtitle:
        tb = add_textbox(slide, title_left, title_top + 1.5, title_w, 0.7)
        write_paragraph(tb.text_frame, subtitle,
                        size=typo.title_size - 4,
                        color=pal.text_secondary, family=typo.family,
                        first=True)

    # Decorative accent lines (fill gap between subtitle and metadata)
    line_y = title_top + 2.5
    add_accent_bar(slide, title_left, line_y, 2.0, 0.04, color=pal.accent)
    add_accent_bar(slide, title_left, line_y + 0.12, 1.2, 0.04,
                   color=pal.accent_soft)

    # Bottom metadata block
    meta_top = layout.slide_height_in - 1.3
    add_line(slide, title_left, meta_top,
             title_left + 6.0, meta_top,
             color=pal.primary, width_pt=1.0)
    meta_y = meta_top + 0.15
    meta_lines = [v for v in [course, student, date] if v]
    for i, line in enumerate(meta_lines):
        tb = add_textbox(slide, title_left, meta_y + i * 0.30, 6.0, 0.28)
        write_paragraph(tb.text_frame, line,
                        size=typo.body_size if i > 0 else typo.body_size + 1,
                        bold=(i == 0), color=pal.text_dark,
                        family=typo.family, first=True)
    return slide
