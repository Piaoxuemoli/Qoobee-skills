"""Comparison slides — pros/cons, before/after, comparison table."""
from __future__ import annotations
from typing import Sequence, Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from ..base import (
    blank_slide, add_chrome, add_rect, add_line, add_textbox,
    write_paragraph, add_oval, add_icon_circle, add_accent_bar,
)
from ..theme import Theme, DEFAULT_THEME


def add_pros_cons(prs, *,
                  title: str,
                  pros: Sequence[str],
                  cons: Sequence[str],
                  pros_label: str = "Pros / Advantages",
                  cons_label: str = "Cons / Disadvantages",
                  page_number=None,
                  theme: Theme = DEFAULT_THEME):
    """Two-column pros/cons with green/red headers and icon circles."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in
    body_bottom = layout.footer_top_in - 0.2
    body_h = body_bottom - body_top
    col_w = (width - 0.4) / 2

    sides = [
        (layout.margin_left_in, pros_label, pros, pal.status_green, "\u2713"),
        (layout.margin_left_in + col_w + 0.4, cons_label, cons,
         pal.status_red, "\u2717"),
    ]
    for left, label, items, color, glyph in sides:
        head_h = 0.5
        # Accent bar at top of header
        add_accent_bar(slide, left, body_top, col_w, 0.05, color=color)
        # Header background
        add_rect(slide, left, body_top + 0.05, col_w, head_h - 0.05,
                 fill=color)
        # Icon circle in header
        add_icon_circle(slide, left + 0.3, body_top + 0.05 + (head_h - 0.05) / 2,
                        0.3, glyph, fill=pal.white, text_color=color,
                        font_size=12, theme=theme)
        # Header text
        tb = add_textbox(slide, left + 0.55, body_top + 0.05,
                         col_w - 0.75, head_h - 0.05,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, label,
                        size=typo.section_title_size + 2, bold=True,
                        color=pal.white, family=typo.family, first=True)

        # Card body
        card_top = body_top + head_h
        card_h = body_h - head_h
        add_rect(slide, left, card_top, col_w, card_h, fill=pal.soft_gray)

        ny = card_top + 0.15
        item_h = max(0.45, (card_h - 0.3) / max(len(items), 1))
        for it in items:
            # Glyph icon
            tb = add_textbox(slide, left + 0.15, ny, 0.35, item_h,
                             anchor=MSO_ANCHOR.TOP)
            write_paragraph(tb.text_frame, glyph,
                            size=typo.body_size + 2, bold=True,
                            color=color, family=typo.family, first=True)
            # Item text
            tb = add_textbox(slide, left + 0.55, ny, col_w - 0.75, item_h,
                             anchor=MSO_ANCHOR.TOP)
            write_paragraph(tb.text_frame, it, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            first=True)
            ny += item_h
    return slide


def add_before_after(prs, *,
                     title: str,
                     left_label: str = "Before",
                     right_label: str = "After",
                     left_items: Sequence[str],
                     right_items: Sequence[str],
                     page_number=None,
                     theme: Theme = DEFAULT_THEME):
    """Before/After comparison with connecting arrow and icon badges."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in
    body_bottom = layout.footer_top_in - 0.2
    body_h = body_bottom - body_top
    arrow_w = 0.9
    col_w = (width - arrow_w) / 2

    sides = [
        (layout.margin_left_in, left_label, left_items, pal.footer_gray, "A"),
        (layout.margin_left_in + col_w + arrow_w, right_label, right_items,
         pal.primary, "B"),
    ]
    for left, label, items, color, badge_text in sides:
        head_h = 0.5
        # Header with accent bar
        add_accent_bar(slide, left, body_top, col_w, 0.05, color=color)
        add_rect(slide, left, body_top + 0.05, col_w, head_h - 0.05,
                 fill=color)
        # Badge in header
        add_icon_circle(slide, left + 0.3, body_top + 0.05 + (head_h - 0.05) / 2,
                        0.3, badge_text, fill=pal.white, text_color=color,
                        font_size=11, theme=theme)
        # Header text
        tb = add_textbox(slide, left + 0.55, body_top + 0.05,
                         col_w - 0.75, head_h - 0.05,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, label,
                        size=typo.section_title_size + 2, bold=True,
                        color=pal.white, family=typo.family, first=True)

        # Card body
        card_top = body_top + head_h
        card_h = body_h - head_h
        add_rect(slide, left, card_top, col_w, card_h, fill=pal.soft_gray)

        tb = add_textbox(slide, left + 0.2, card_top + 0.15,
                         col_w - 0.4, card_h - 0.3)
        first = True
        for it in items:
            write_paragraph(tb.text_frame, it, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=4, first=first)
            first = False

    # Arrow between columns
    arr_x = layout.margin_left_in + col_w + 0.1
    arr_y = body_top + body_h / 2 - 0.4
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(arr_x), Inches(arr_y),
                                 Inches(arrow_w - 0.2), Inches(0.8))
    arr.shadow.inherit = False
    arr.fill.solid()
    arr.fill.fore_color.rgb = pal.primary
    arr.line.fill.background()
    return slide
