"""Timeline and process flow slides."""
from __future__ import annotations
from typing import Sequence, Optional, Dict

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from ..base import (
    blank_slide, add_chrome, add_rect, add_line, add_textbox,
    write_paragraph, add_oval,
)
from ..theme import Theme, DEFAULT_THEME


def add_process_flow(prs, *,
                     title: str,
                     steps: Sequence[str],
                     page_number=None,
                     theme: Theme = DEFAULT_THEME):
    """Horizontal process flow with arrows between steps."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in + 0.5
    n = max(len(steps), 1)
    arrow_space = 0.5 * (n - 1)
    step_w = (width - arrow_space) / n
    step_h = 2.0

    for i, step in enumerate(steps):
        x = layout.margin_left_in + i * (step_w + 0.5)
        # Step box
        add_rect(slide, x, body_top, step_w, step_h,
                 fill=pal.primary if i == 0 else pal.soft_gray,
                 line=pal.primary, line_width=1.0)
        # Step number circle
        num_d = 0.5
        cx = x + step_w / 2 - num_d / 2
        cy = body_top - num_d / 2
        add_oval(slide, cx, cy, num_d, num_d, fill=pal.primary)
        tb = add_textbox(slide, cx, cy, num_d, num_d,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, str(i + 1),
                        size=typo.body_size, bold=True,
                        color=pal.white, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)
        # Step text
        tb = add_textbox(slide, x + 0.15, body_top + 0.3,
                         step_w - 0.3, step_h - 0.4)
        write_paragraph(tb.text_frame, step, size=typo.body_size,
                        color=pal.text_dark if i > 0 else pal.white,
                        family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)
        # Arrow between steps
        if i < n - 1:
            arr_x = x + step_w + 0.05
            arr_y = body_top + step_h / 2 - 0.2
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         Inches(arr_x), Inches(arr_y),
                                         Inches(0.4), Inches(0.4))
            arr.shadow.inherit = False
            arr.fill.solid()
            arr.fill.fore_color.rgb = pal.primary
            arr.line.fill.background()
    return slide


def add_phases(prs, *,
               title: str,
               phases: Sequence[Dict],
               page_number=None,
               theme: Theme = DEFAULT_THEME):
    """3-phase chevron timeline.
    phases: [{"label": str, "timeframe": str, "deliverables": [str]}]
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in
    n = max(len(phases), 1)
    chev_w = (width + 0.4) / n - 0.1
    chev_h = 0.6
    chev_y = body_top + 0.4

    for i, ph in enumerate(phases):
        x = layout.margin_left_in + i * (chev_w - 0.4)
        fill = pal.primary if i == 0 else pal.primary_light

        # Chevron shape
        s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                   Inches(x), Inches(chev_y),
                                   Inches(chev_w), Inches(chev_h))
        s.shadow.inherit = False
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.fill.background()

        # Label inside chevron
        tb = add_textbox(slide, x + 0.3, chev_y, chev_w - 0.6, chev_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, ph.get("label", f"Phase {i+1}"),
                        size=typo.body_size, bold=True,
                        color=pal.white, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)

        # Timeframe below
        tf = ph.get("timeframe", "")
        if tf:
            tb = add_textbox(slide, x + 0.2, chev_y + chev_h + 0.1,
                             chev_w - 0.4, 0.3)
            write_paragraph(tb.text_frame, tf,
                            size=typo.small_size, color=pal.text_secondary,
                            family=typo.family,
                            align=PP_ALIGN.CENTER, first=True)

        # Deliverables below timeframe
        deliv_top = chev_y + chev_h + 0.5
        tb = add_textbox(slide, x + 0.2, deliv_top,
                         chev_w - 0.4,
                         layout.footer_top_in - deliv_top - 0.2)
        first = True
        for d in ph.get("deliverables", []):
            write_paragraph(tb.text_frame, d, size=typo.small_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=2, first=first)
            first = False
    return slide
