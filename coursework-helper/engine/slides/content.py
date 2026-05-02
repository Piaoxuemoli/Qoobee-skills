"""Content slides — bullet list, two-column, question card."""
from __future__ import annotations
from typing import Sequence, Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..base import (
    blank_slide, add_chrome, add_rect, add_line, add_textbox,
    write_paragraph, add_oval, add_footer,
)
from ..theme import Theme, DEFAULT_THEME


def add_bullet_list(prs, *,
                    title: str,
                    bullets: Sequence[str],
                    key_message: Optional[str] = None,
                    page_number=None,
                    theme: Theme = DEFAULT_THEME):
    """Standard content slide with title + bullet points."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in

    if key_message:
        tb = add_textbox(slide, layout.margin_left_in, body_top, width, 0.4)
        write_paragraph(tb.text_frame, key_message,
                        size=typo.body_size + 1, bold=True,
                        color=pal.primary, family=typo.family, first=True)
        body_top += 0.5

    tb = add_textbox(slide, layout.margin_left_in, body_top, width,
                     layout.footer_top_in - body_top - 0.2)
    first = True
    for b in bullets:
        write_paragraph(tb.text_frame, b, size=typo.body_size,
                        color=pal.text_dark, family=typo.family,
                        bullet=True, space_after=6, first=first)
        first = False
    return slide


def add_two_column(prs, *,
                   title: str,
                   left_title: str,
                   left_bullets: Sequence[str],
                   right_title: str,
                   right_bullets: Sequence[str],
                   key_message: Optional[str] = None,
                   page_number=None,
                   theme: Theme = DEFAULT_THEME):
    """Two-column content slide."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in
    col_gap = 0.4
    col_w = (width - col_gap) / 2

    for ci, (col_title, col_bullets, col_left) in enumerate([
        (left_title, left_bullets, layout.margin_left_in),
        (right_title, right_bullets, layout.margin_left_in + col_w + col_gap),
    ]):
        # Column header
        add_rect(slide, col_left, body_top, col_w, 0.45, fill=pal.primary)
        tb = add_textbox(slide, col_left + 0.15, body_top, col_w - 0.3, 0.45,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, col_title,
                        size=typo.section_title_size, bold=True,
                        color=pal.white, family=typo.family, first=True)
        # Column body
        card_top = body_top + 0.45
        card_h = layout.footer_top_in - card_top - 0.2
        add_rect(slide, col_left, card_top, col_w, card_h, fill=pal.soft_gray)
        tb = add_textbox(slide, col_left + 0.2, card_top + 0.15,
                         col_w - 0.4, card_h - 0.3)
        first = True
        for b in col_bullets:
            write_paragraph(tb.text_frame, b, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=4, first=first)
            first = False
    return slide


def add_question(prs, *,
                 question: str,
                 subtitle: Optional[str] = None,
                 page_number=None,
                 theme: Theme = DEFAULT_THEME):
    """Hook/question slide — a single provocative question on a clean background."""
    slide = blank_slide(prs)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Light background tint
    add_rect(slide, 0, 0, layout.slide_width_in, layout.slide_height_in,
             fill=pal.soft_gray)

    # Big question mark accent
    tb = add_textbox(slide, layout.margin_left_in + 0.3, 1.5, 1.5, 2.0)
    write_paragraph(tb.text_frame, "?",
                    size=120, bold=True, color=pal.primary_light,
                    family=typo.family, first=True)

    # Question text
    q_left = layout.margin_left_in + 2.0
    q_w = layout.slide_width_in - q_left - layout.margin_right_in - 0.5
    tb = add_textbox(slide, q_left, 2.2, q_w, 2.0)
    write_paragraph(tb.text_frame, question,
                    size=typo.title_size + 4, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)

    if subtitle:
        tb = add_textbox(slide, q_left, 4.3, q_w, 0.8)
        write_paragraph(tb.text_frame, subtitle,
                        size=typo.body_size + 2,
                        color=pal.text_secondary, family=typo.family,
                        first=True)

    add_footer(slide, theme, page_number=page_number)
    return slide
