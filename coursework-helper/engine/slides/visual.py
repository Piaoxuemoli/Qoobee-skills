"""Visual slides — diagram placeholder, simple data display."""
from __future__ import annotations
from typing import Sequence, Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_chrome, add_rect, add_line, add_textbox,
    write_paragraph, add_oval, add_footer,
)
from ..theme import Theme, DEFAULT_THEME


def add_diagram(prs, *,
                title: str,
                items: Sequence[str],
                key_message: Optional[str] = None,
                page_number=None,
                theme: Theme = DEFAULT_THEME):
    """Diagram placeholder slide — items arranged in a grid with connecting lines.
    Good for frameworks, models, conceptual diagrams.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in

    if key_message:
        tb = add_textbox(slide, layout.margin_left_in, body_top, width, 0.35)
        write_paragraph(tb.text_frame, key_message,
                        size=typo.body_size, bold=True,
                        color=pal.primary, family=typo.family, first=True)
        body_top += 0.5

    # Arrange items in a grid (2 or 3 columns)
    n = min(len(items), 6)
    if n <= 2:
        cols = n
    elif n <= 4:
        cols = 2
    else:
        cols = 3
    rows = (n + cols - 1) // cols

    gap = 0.3
    card_w = (width - gap * (cols - 1)) / cols
    avail_h = layout.footer_top_in - body_top - 0.2
    card_h = (avail_h - gap * (rows - 1)) / rows

    for i, item in enumerate(items[:6]):
        r, c = divmod(i, cols)
        x = layout.margin_left_in + c * (card_w + gap)
        y = body_top + r * (card_h + gap)
        add_rect(slide, x, y, card_w, card_h,
                 fill=pal.soft_gray, line=pal.primary, line_width=1.0)
        tb = add_textbox(slide, x + 0.15, y + 0.15,
                         card_w - 0.3, card_h - 0.3,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, item,
                        size=typo.body_size, color=pal.text_dark,
                        family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)
    return slide


def add_table_slide(prs, *,
                    title: str,
                    headers: Sequence[str],
                    rows: Sequence[Sequence[str]],
                    page_number=None,
                    theme: Theme = DEFAULT_THEME):
    """Simple table slide using native PPTX table."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in + 0.2
    avail_h = layout.footer_top_in - body_top - 0.3

    n_cols = max(len(headers), 1)
    n_rows_data = len(rows)
    n_total_rows = n_rows_data + 1  # +1 for header

    table = slide.shapes.add_table(
        n_total_rows, n_cols,
        Inches(layout.margin_left_in), Inches(body_top),
        Inches(width), Inches(min(avail_h, n_total_rows * 0.55))
    ).table

    # Style header
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(typo.body_size)
                run.font.bold = True
                run.font.color.rgb = pal.white
                run.font.name = typo.family
        cell.fill.solid()
        cell.fill.fore_color.rgb = pal.primary

    # Style body rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(typo.body_size)
                    run.font.color.rgb = pal.text_dark
                    run.font.name = typo.family
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = pal.soft_gray

    return slide
