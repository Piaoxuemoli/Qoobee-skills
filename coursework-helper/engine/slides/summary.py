"""Summary slides — three takeaways, closing, quote."""
from __future__ import annotations
from typing import Sequence, Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..base import (
    blank_slide, add_chrome, add_rect, add_line, add_textbox,
    write_paragraph, add_oval, add_footer, add_icon_circle, add_accent_bar,
)
from ..theme import Theme, DEFAULT_THEME


def add_three_takeaways(prs, *,
                        title: str = "Summary",
                        takeaways: Sequence[str],
                        page_number=None,
                        theme: Theme = DEFAULT_THEME):
    """Key takeaway cards with numbered circles and accent bars.

    Card height adapts to content: capped at 2.5" per card when there are
    few takeaways, so the slide doesn't look empty.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in + 0.2
    n = min(len(takeaways), 4)
    card_gap = 0.25
    card_w = (width - card_gap * (n - 1)) / n

    # Adaptive card height: use content-based height, capped
    max_card_h = layout.footer_top_in - body_top - 0.3
    card_h = min(max_card_h, 2.5)  # don't stretch cards beyond 2.5"

    for i, tk in enumerate(takeaways[:4]):
        x = layout.margin_left_in + i * (card_w + card_gap)

        # Accent bar at top of card
        add_accent_bar(slide, x, body_top, card_w, 0.06, color=pal.accent)

        # Card background
        add_rect(slide, x, body_top + 0.06, card_w, card_h - 0.06,
                 fill=pal.soft_gray, line=pal.primary, line_width=0.75)

        # Number circle
        num_d = 0.5
        add_icon_circle(slide, x + 0.3, body_top + 0.4, num_d,
                        str(i + 1), fill=pal.primary, theme=theme)

        # Takeaway text — left-aligned for readability
        tb = add_textbox(slide, x + 0.15, body_top + 0.8,
                         card_w - 0.3, card_h - 1.0)
        write_paragraph(tb.text_frame, tk, size=typo.body_size,
                        color=pal.text_dark, family=typo.family,
                        align=PP_ALIGN.LEFT, first=True)
    return slide


def add_closing(prs, *,
                title: str = "Thank You / Q&A",
                message: Optional[str] = None,
                page_number=None,
                theme: Theme = DEFAULT_THEME):
    """Closing/Q&A slide with centered text on accent background."""
    slide = blank_slide(prs)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Full-bleed primary color background
    add_rect(slide, 0, 0, layout.slide_width_in, layout.slide_height_in,
             fill=pal.primary)

    # Decorative top rule
    add_line(slide, 3.0, 1.8, layout.slide_width_in - 3.0, 1.8,
             color=pal.white, width_pt=1.0)

    # Title centered
    tb = add_textbox(slide, layout.margin_left_in, 2.2,
                     layout.slide_width_in - layout.margin_left_in - layout.margin_right_in,
                     1.5, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, title,
                    size=typo.title_size + 20, bold=True,
                    color=pal.white, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)

    if message:
        tb = add_textbox(slide, 2.0, 4.0,
                         layout.slide_width_in - 4.0, 1.0)
        write_paragraph(tb.text_frame, message,
                        size=typo.body_size + 2,
                        color=pal.white, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)

    # Decorative bottom rule
    add_line(slide, 3.0, 5.5, layout.slide_width_in - 3.0, 5.5,
             color=pal.white, width_pt=0.5)
    return slide


def add_quote(prs, *,
              quote: str,
              author: str,
              title: Optional[str] = None,
              page_number=None,
              theme: Theme = DEFAULT_THEME):
    """Pull-quote slide with large quotation marks and accent bar."""
    slide = blank_slide(prs)
    if title:
        add_chrome(slide, title=title, theme=theme, page_number=page_number)
        body_top = 2.0
    else:
        body_top = 1.5
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Left accent bar
    add_accent_bar(slide, 0.5, body_top, 0.08, 4.0, color=pal.accent)

    # Big opening quote mark
    tb = add_textbox(slide, 0.8, body_top, 1.5, 1.5)
    write_paragraph(tb.text_frame, "\u201C",
                    size=typo.title_size + 60, bold=True,
                    color=pal.primary, family=typo.family, first=True)

    # Quote body
    q_left = 2.0
    q_w = layout.slide_width_in - q_left - layout.margin_right_in - 0.5
    tb = add_textbox(slide, q_left, body_top + 0.4, q_w, 3.5)
    write_paragraph(tb.text_frame, quote,
                    size=typo.title_size, color=pal.text_dark,
                    family=typo.family, first=True)

    # Attribution
    attr_y = body_top + 4.0
    add_line(slide, q_left, attr_y, q_left + 0.6, attr_y,
             color=pal.primary, width_pt=2.0)
    tb = add_textbox(slide, q_left, attr_y + 0.1, q_w, 0.35)
    write_paragraph(tb.text_frame, author,
                    size=typo.body_size + 2, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)
    return slide
