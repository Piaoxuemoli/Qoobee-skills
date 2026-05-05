"""Extract text and images from PPT/PDF/DOCX files.

Outputs per-source directories with text files and extracted images,
plus a manifest JSON summarizing what was extracted.

Usage:
    python extract_content.py \
        --inputs "file1.pptx|file2.pdf|file3.docx" \
        --output "outputs/course/01_extracted/" \
        --manifest "outputs/course/00_admin/extract_manifest.json"
"""
from __future__ import annotations
import argparse
import json
import os
import re
import sys
from pathlib import Path
from typing import Dict, List, Any


def _safe_name(name: str) -> str:
    """Convert a filename to a safe directory name."""
    name = Path(name).stem
    # Replace problematic characters
    name = re.sub(r'[<>:"/\\|?*\s]+', '_', name)
    return name[:64] or "unnamed"


def _extract_pptx(file_path: str, out_dir: Path) -> Dict[str, Any]:
    """Extract text and images from a PPTX file."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    text_count = 0
    img_count = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        img_in_slide = 0

        for shape in slide.shapes:
            # Extract text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text_parts.append(text)

            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_in_slide += 1
                img_name = f"slide_{slide_idx:02d}_img_{img_in_slide:02d}.png"
                img_path = out_dir / img_name
                with open(img_path, "wb") as f:
                    f.write(shape.image.blob)
                img_count += 1

            # Extract images from groups
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_in_slide += 1
                        img_name = f"slide_{slide_idx:02d}_img_{img_in_slide:02d}.png"
                        img_path = out_dir / img_name
                        with open(img_path, "wb") as f:
                            f.write(sub_shape.image.blob)
                        img_count += 1

        # Write slide text
        if slide_text_parts:
            txt_name = f"slide_{slide_idx:02d}.txt"
            (out_dir / txt_name).write_text(
                "\n".join(slide_text_parts), encoding="utf-8")
            text_count += 1

    return {"text_files": text_count, "image_files": img_count,
            "slides": len(prs.slides)}


def _extract_pdf(file_path: str, out_dir: Path) -> Dict[str, Any]:
    """Extract text and images from a PDF file using PyMuPDF."""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    text_count = 0
    img_count = 0
    page_count = len(doc)

    for page_idx, page in enumerate(doc, start=1):
        # Extract text
        text = page.get_text("text").strip()
        if text:
            txt_name = f"page_{page_idx:02d}.txt"
            (out_dir / txt_name).write_text(text, encoding="utf-8")
            text_count += 1

        # Extract images
        img_list = page.get_images(full=True)
        img_in_page = 0
        for img_info in img_list:
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                if base_image and base_image.get("image"):
                    img_in_page += 1
                    ext = base_image.get("ext", "png")
                    img_name = f"page_{page_idx:02d}_img_{img_in_page:02d}.{ext}"
                    img_path = out_dir / img_name
                    with open(img_path, "wb") as f:
                        f.write(base_image["image"])
                    img_count += 1
            except Exception:
                pass  # Skip unextractable images

    doc.close()
    return {"text_files": text_count, "image_files": img_count,
            "pages": page_count}


def _extract_docx(file_path: str, out_dir: Path) -> Dict[str, Any]:
    """Extract text and images from a DOCX file."""
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    doc = Document(file_path)
    text_count = 0
    img_count = 0

    # Extract text by paragraphs
    text_parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            text_parts.append(text)

    if text_parts:
        (out_dir / "content.txt").write_text(
            "\n".join(text_parts), encoding="utf-8")
        text_count = 1

    # Extract inline images
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            img_count += 1
            img_name = f"image_{img_count:02d}.png"
            img_path = out_dir / img_name
            with open(img_path, "wb") as f:
                f.write(rel.target_part.blob)

    return {"text_files": text_count, "image_files": img_count}


def _extract_text(file_path: str, out_dir: Path) -> Dict[str, Any]:
    """Copy text/markdown files directly."""
    content = Path(file_path).read_text(encoding="utf-8", errors="replace")
    ext = Path(file_path).suffix
    out_name = f"content{ext}" if ext else "content.txt"
    (out_dir / out_name).write_text(content, encoding="utf-8")
    return {"text_files": 1, "image_files": 0}


_EXTRACTORS = {
    ".pptx": _extract_pptx,
    ".ppt": _extract_pptx,
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".doc": _extract_docx,
    ".md": _extract_text,
    ".txt": _extract_text,
    ".rst": _extract_text,
    ".tex": _extract_text,
}


def extract_file(file_path: str, output_base: Path) -> Dict[str, Any]:
    """Extract content from a single file. Returns source info dict."""
    path = Path(file_path)
    ext = path.suffix.lower()
    extractor = _EXTRACTORS.get(ext)

    if extractor is None:
        return {
            "name": path.stem,
            "type": ext.lstrip("."),
            "path": str(path),
            "status": "skipped",
            "reason": f"Unsupported format: {ext}",
        }

    safe = _safe_name(path.name)
    out_dir = output_base / safe
    out_dir.mkdir(parents=True, exist_ok=True)

    try:
        stats = extractor(str(path), out_dir)
        return {
            "name": safe,
            "type": ext.lstrip("."),
            "path": str(out_dir.relative_to(output_base.parent.parent)),
            "status": "ok",
            **stats,
        }
    except Exception as e:
        return {
            "name": safe,
            "type": ext.lstrip("."),
            "path": str(path),
            "status": "error",
            "reason": str(e),
        }


def main():
    parser = argparse.ArgumentParser(
        description="Extract text and images from course materials")
    parser.add_argument("--inputs", required=True,
                        help="Pipe-separated list of input files or directories")
    parser.add_argument("--output", required=True,
                        help="Output directory for extracted content")
    parser.add_argument("--manifest", required=True,
                        help="Path to write extract_manifest.json")
    args = parser.parse_args()

    output_dir = Path(args.output)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Collect all input files
    input_paths = [p.strip() for p in args.inputs.split("|") if p.strip()]
    files_to_process: List[str] = []

    for p in input_paths:
        p_path = Path(p)
        if p_path.is_dir():
            # Walk directory, skip hidden/system dirs
            for f in sorted(p_path.rglob("*")):
                if f.is_file() and f.suffix.lower() in _EXTRACTORS:
                    files_to_process.append(str(f))
        elif p_path.is_file():
            files_to_process.append(str(p))
        else:
            print(f"Warning: not found: {p}", file=sys.stderr)

    if not files_to_process:
        print("Error: no processable files found", file=sys.stderr)
        sys.exit(1)

    # Extract each file
    sources = []
    total_text = 0
    total_images = 0

    for fpath in files_to_process:
        print(f"Extracting: {Path(fpath).name}")
        info = extract_file(fpath, output_dir)
        sources.append(info)
        if info.get("status") == "ok":
            total_text += info.get("text_files", 0)
            total_images += info.get("image_files", 0)

    # Write manifest
    manifest = {
        "sources": sources,
        "totals": {
            "source_count": len(sources),
            "text_files": total_text,
            "image_files": total_images,
            "ok": sum(1 for s in sources if s.get("status") == "ok"),
            "errors": sum(1 for s in sources if s.get("status") == "error"),
            "skipped": sum(1 for s in sources if s.get("status") == "skipped"),
        },
    }

    manifest_path = Path(args.manifest)
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"\nDone: {total_text} text files, {total_images} images "
          f"from {len(sources)} sources")
    print(f"Manifest: {manifest_path}")


if __name__ == "__main__":
    main()
