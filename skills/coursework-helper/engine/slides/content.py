"""Content slides — bullet list, two-column, question card."""
from __future__ import annotations
from typing import Sequence, Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..base import (
    blank_slide, add_chrome, add_rect, add_line, add_textbox,
    write_paragraph, add_oval, add_footer, add_icon_circle, add_accent_bar,
)
from ..theme import Theme, DEFAULT_THEME


def add_bullet_list(prs, *,
                    title: str,
                    bullets: Sequence[str],
                    key_message: Optional[str] = None,
                    page_number=None,
                    theme: Theme = DEFAULT_THEME):
    """Standard content slide with title + bullet points.

    When there are 3 or fewer bullets, uses icon-circle + card layout for
    visual density. 4+ bullets use compact list with tighter spacing.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in

    if key_message:
        tb = add_textbox(slide, layout.margin_left_in, body_top, width, 0.4)
        write_paragraph(tb.text_frame, key_message,
                        size=typo.body_size + 1, bold=True,
                        color=pal.primary, family=typo.family, first=True)
        body_top += 0.5

    if len(bullets) <= 3:
        # Icon-card layout: each bullet gets its own card with icon circle
        avail_h = layout.footer_top_in - body_top - 0.3
        card_h = avail_h / len(bullets)
        card_h = min(card_h, 1.6)  # cap card height
        gap = 0.15

        for i, b in enumerate(bullets):
            cy = body_top + i * (card_h + gap)
            # Accent bar on left edge of card
            add_accent_bar(slide, layout.margin_left_in, cy, 0.06, card_h,
                           color=pal.accent)
            # Card background
            add_rect(slide, layout.margin_left_in + 0.06, cy,
                     width - 0.06, card_h, fill=pal.soft_gray)
            # Icon circle
            icon_d = min(0.45, card_h - 0.2)
            add_icon_circle(slide,
                            layout.margin_left_in + 0.06 + 0.35,
                            cy + card_h / 2,
                            icon_d, str(i + 1),
                            fill=pal.primary, theme=theme)
            # Text
            tb = add_textbox(slide, layout.margin_left_in + 0.95, cy + 0.1,
                             width - 1.1, card_h - 0.2,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, b, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            first=True)
    else:
        # Compact list: tighter spacing for 4+ bullets
        tb = add_textbox(slide, layout.margin_left_in, body_top, width,
                         layout.footer_top_in - body_top - 0.2)
        first = True
        for b in bullets:
            write_paragraph(tb.text_frame, b, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=4, first=first)
            first = False
    return slide


def add_two_column(prs, *,
                   title: str,
                   left_title: str,
                   left_bullets: Sequence[str],
                   right_title: str,
                   right_bullets: Sequence[str],
                   key_message: Optional[str] = None,
                   page_number=None,
                   theme: Theme = DEFAULT_THEME):
    """Two-column content slide with colored headers and icon badges."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in
    col_gap = 0.4
    col_w = (width - col_gap) / 2

    for ci, (col_title, col_bullets, col_left) in enumerate([
        (left_title, left_bullets, layout.margin_left_in),
        (right_title, right_bullets, layout.margin_left_in + col_w + col_gap),
    ]):
        # Column header with accent
        add_rect(slide, col_left, body_top, col_w, 0.5, fill=pal.primary)
        # Icon circle in header
        add_icon_circle(slide, col_left + 0.3, body_top + 0.25, 0.32,
                        str(ci + 1), fill=pal.white, text_color=pal.primary,
                        font_size=12, theme=theme)
        # Header text
        tb = add_textbox(slide, col_left + 0.55, body_top, col_w - 0.7, 0.5,
                         anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb.text_frame, col_title,
                        size=typo.section_title_size, bold=True,
                        color=pal.white, family=typo.family, first=True)
        # Column body card
        card_top = body_top + 0.5
        card_h = layout.footer_top_in - card_top - 0.2
        add_rect(slide, col_left, card_top, col_w, card_h, fill=pal.soft_gray)
        tb = add_textbox(slide, col_left + 0.2, card_top + 0.15,
                         col_w - 0.4, card_h - 0.3)
        first = True
        for b in col_bullets:
            write_paragraph(tb.text_frame, b, size=typo.body_size,
                            color=pal.text_dark, family=typo.family,
                            bullet=True, space_after=4, first=first)
            first = False
    return slide


def add_question(prs, *,
                 question: str,
                 subtitle: Optional[str] = None,
                 page_number=None,
                 theme: Theme = DEFAULT_THEME):
    """Hook/question slide — a single provocative question on a clean background."""
    slide = blank_slide(prs)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    # Light background tint
    add_rect(slide, 0, 0, layout.slide_width_in, layout.slide_height_in,
             fill=pal.soft_gray)

    # Big question mark accent
    tb = add_textbox(slide, layout.margin_left_in + 0.3, 1.5, 1.5, 2.0)
    write_paragraph(tb.text_frame, "?",
                    size=120, bold=True, color=pal.primary_light,
                    family=typo.family, first=True)

    # Question text
    q_left = layout.margin_left_in + 2.0
    q_w = layout.slide_width_in - q_left - layout.margin_right_in - 0.5
    tb = add_textbox(slide, q_left, 2.2, q_w, 2.0)
    write_paragraph(tb.text_frame, question,
                    size=typo.title_size + 4, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)

    if subtitle:
        tb = add_textbox(slide, q_left, 4.3, q_w, 0.8)
        write_paragraph(tb.text_frame, subtitle,
                        size=typo.body_size + 2,
                        color=pal.text_secondary, family=typo.family,
                        first=True)

    add_footer(slide, theme, page_number=page_number)
    return slide
