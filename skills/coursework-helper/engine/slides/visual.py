"""Visual slides — diagram, icon grid, simple data display."""
from __future__ import annotations
from typing import Sequence, Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from ..base import (
    blank_slide, add_chrome, add_rect, add_line, add_textbox,
    write_paragraph, add_oval, add_footer, add_icon_circle, add_accent_bar,
)
from ..theme import Theme, DEFAULT_THEME


def add_diagram(prs, *,
                title: str,
                items: Sequence[str],
                key_message: Optional[str] = None,
                page_number=None,
                theme: Theme = DEFAULT_THEME):
    """Diagram slide — items arranged in a grid with icon circles and accent bars.

    For 1-2 items: large cards with icon circle + text + accent bar.
    For 3-6 items: grid with numbered cards.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in

    if key_message:
        tb = add_textbox(slide, layout.margin_left_in, body_top, width, 0.35)
        write_paragraph(tb.text_frame, key_message,
                        size=typo.body_size, bold=True,
                        color=pal.primary, family=typo.family, first=True)
        body_top += 0.5

    n = min(len(items), 6)
    avail_h = layout.footer_top_in - body_top - 0.2

    if n <= 2:
        # Large card layout: each item gets a full-width card with icon + text
        card_h = min(avail_h / n - 0.15, 2.2)
        for i, item in enumerate(items[:2]):
            y = body_top + i * (card_h + 0.15)
            # Accent bar on left
            add_accent_bar(slide, layout.margin_left_in, y, 0.06, card_h,
                           color=pal.accent)
            # Card background
            add_rect(slide, layout.margin_left_in + 0.06, y,
                     width - 0.06, card_h, fill=pal.soft_gray)
            # Icon circle
            icon_d = min(0.55, card_h - 0.3)
            add_icon_circle(slide,
                            layout.margin_left_in + 0.5, y + card_h / 2,
                            icon_d, str(i + 1),
                            fill=pal.primary, theme=theme)
            # Text
            tb = add_textbox(slide, layout.margin_left_in + 1.0, y + 0.15,
                             width - 1.2, card_h - 0.3,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, item, size=typo.body_size + 1,
                            color=pal.text_dark, family=typo.family,
                            first=True)
    else:
        # Grid layout with numbered icon circles
        if n <= 4:
            cols = 2
        else:
            cols = 3
        rows = (n + cols - 1) // cols

        gap = 0.25
        card_w = (width - gap * (cols - 1)) / cols
        card_h = (avail_h - gap * (rows - 1)) / rows

        for i, item in enumerate(items[:6]):
            r, c = divmod(i, cols)
            x = layout.margin_left_in + c * (card_w + gap)
            y = body_top + r * (card_h + gap)

            # Accent bar at top
            add_accent_bar(slide, x, y, card_w, 0.05, color=pal.accent)
            # Card background
            add_rect(slide, x, y + 0.05, card_w, card_h - 0.05,
                     fill=pal.soft_gray, line=pal.primary, line_width=0.75)
            # Icon circle
            icon_d = 0.4
            add_icon_circle(slide, x + card_w / 2, y + 0.35,
                            icon_d, str(i + 1),
                            fill=pal.primary, theme=theme)
            # Text
            tb = add_textbox(slide, x + 0.1, y + 0.65,
                             card_w - 0.2, card_h - 0.8,
                             anchor=MSO_ANCHOR.MIDDLE)
            write_paragraph(tb.text_frame, item,
                            size=typo.body_size, color=pal.text_dark,
                            family=typo.family,
                            align=PP_ALIGN.CENTER, first=True)
    return slide


def add_icon_grid(prs, *,
                  title: str,
                  items: Sequence[dict],
                  key_message: Optional[str] = None,
                  columns: int = 3,
                  page_number=None,
                  theme: Theme = DEFAULT_THEME):
    """Icon grid slide — N cards with icon circle + title + description.

    Each item dict should have: {"title": str, "desc": str, "icon": str (optional)}
    Icons default to sequential numbers if not provided.
    """
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in

    if key_message:
        tb = add_textbox(slide, layout.margin_left_in, body_top, width, 0.35)
        write_paragraph(tb.text_frame, key_message,
                        size=typo.body_size, bold=True,
                        color=pal.primary, family=typo.family, first=True)
        body_top += 0.5

    n = min(len(items), 6)
    cols = min(columns, n)
    rows = (n + cols - 1) // cols
    gap = 0.25
    card_w = (width - gap * (cols - 1)) / cols
    avail_h = layout.footer_top_in - body_top - 0.2
    card_h = (avail_h - gap * (rows - 1)) / rows

    for i, item in enumerate(items[:6]):
        r, c = divmod(i, cols)
        x = layout.margin_left_in + c * (card_w + gap)
        y = body_top + r * (card_h + gap)

        icon_text = item.get("icon", str(i + 1))
        item_title = item.get("title", "")
        item_desc = item.get("desc", "")

        # Accent bar at top
        add_accent_bar(slide, x, y, card_w, 0.05, color=pal.accent)
        # Card background
        add_rect(slide, x, y + 0.05, card_w, card_h - 0.05,
                 fill=pal.soft_gray, line=pal.primary, line_width=0.75)
        # Icon circle
        icon_d = 0.45
        add_icon_circle(slide, x + card_w / 2, y + 0.4, icon_d,
                        icon_text, fill=pal.primary, theme=theme)
        # Title
        title_y = y + 0.75
        tb = add_textbox(slide, x + 0.1, title_y, card_w - 0.2, 0.35)
        write_paragraph(tb.text_frame, item_title,
                        size=typo.body_size, bold=True,
                        color=pal.text_dark, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)
        # Description
        if item_desc:
            desc_y = title_y + 0.35
            tb = add_textbox(slide, x + 0.1, desc_y,
                             card_w - 0.2, card_h - 1.2)
            write_paragraph(tb.text_frame, item_desc,
                            size=typo.small_size,
                            color=pal.text_secondary, family=typo.family,
                            align=PP_ALIGN.CENTER, first=True)
    return slide


def add_table_slide(prs, *,
                    title: str,
                    headers: Sequence[str],
                    rows: Sequence[Sequence[str]],
                    page_number=None,
                    theme: Theme = DEFAULT_THEME):
    """Simple table slide using native PPTX table."""
    slide = blank_slide(prs)
    add_chrome(slide, title=title, theme=theme, page_number=page_number)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in
    body_top = layout.body_top_in + 0.2
    avail_h = layout.footer_top_in - body_top - 0.3

    n_cols = max(len(headers), 1)
    n_rows_data = len(rows)
    n_total_rows = n_rows_data + 1  # +1 for header

    table = slide.shapes.add_table(
        n_total_rows, n_cols,
        Inches(layout.margin_left_in), Inches(body_top),
        Inches(width), Inches(min(avail_h, n_total_rows * 0.55))
    ).table

    # Style header
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(typo.body_size)
                run.font.bold = True
                run.font.color.rgb = pal.white
                run.font.name = typo.family
        cell.fill.solid()
        cell.fill.fore_color.rgb = pal.primary

    # Style body rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(typo.body_size)
                    run.font.color.rgb = pal.text_dark
                    run.font.name = typo.family
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = pal.soft_gray

    return slide
