"""Section divider and stat-hero slides."""
from __future__ import annotations
from typing import Optional

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..base import (
    blank_slide, add_rect, add_line, add_textbox, write_paragraph,
    enable_text_shrink, add_oval, add_title, add_footer, add_accent_bar,
)
from ..theme import Theme, DEFAULT_THEME


def add_section_divider(prs, *,
                        section_number: str,
                        section_title: str,
                        subtitle: Optional[str] = None,
                        page_number=None,
                        theme: Theme = DEFAULT_THEME):
    """Full-bleed section divider with left panel and decorative lines."""
    slide = blank_slide(prs)
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    panel_w = 4.5
    add_rect(slide, 0, 0, panel_w, layout.slide_height_in, fill=pal.primary)

    # Section number
    tb = add_textbox(slide, 0.5, layout.slide_height_in / 2 - 1.5,
                     panel_w - 1.0, 2.0, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, str(section_number),
                    size=typo.title_size + 56, bold=True,
                    color=pal.white, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)

    # Right side title
    right_left = panel_w + 0.6
    right_w = layout.slide_width_in - right_left - layout.margin_right_in
    tb = add_textbox(slide, right_left,
                     layout.slide_height_in / 2 - 1.0,
                     right_w, 1.4, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, section_title,
                    size=typo.title_size + 8, bold=True,
                    color=pal.text_dark, family=typo.family, first=True)

    # Accent underline
    add_line(slide, right_left, layout.slide_height_in / 2 + 0.40,
             right_left + 1.6, layout.slide_height_in / 2 + 0.40,
             color=pal.primary, width_pt=2.5)

    # Decorative parallel lines below the title
    line_y = layout.slide_height_in / 2 + 0.65
    for offset in [0, 0.12, 0.24]:
        add_line(slide, right_left, line_y + offset,
                 right_left + 0.8, line_y + offset,
                 color=pal.accent_soft, width_pt=0.5)

    if subtitle:
        tb = add_textbox(slide, right_left,
                         layout.slide_height_in / 2 + 1.1, right_w, 1.0)
        write_paragraph(tb.text_frame, subtitle,
                        size=typo.body_size + 2, color=pal.footer_gray,
                        family=typo.family, first=True)
    return slide


def add_stat_hero(prs, *,
                  title: Optional[str] = None,
                  stat: str,
                  stat_label: str,
                  context: Optional[str] = None,
                  page_number=None,
                  theme: Theme = DEFAULT_THEME):
    """One big number + label + optional context, with accent decorations."""
    slide = blank_slide(prs)
    if title:
        add_chrome_title(slide, title, theme, page_number)
        body_top = 2.0
    else:
        body_top = 1.0
    pal, typo, layout = theme.palette, theme.typography, theme.layout

    width = layout.slide_width_in - layout.margin_left_in - layout.margin_right_in

    # Big number
    num_h = 2.2
    tb = add_textbox(slide, layout.margin_left_in, body_top + 0.3,
                     width, num_h, anchor=MSO_ANCHOR.MIDDLE)
    write_paragraph(tb.text_frame, stat,
                    size=typo.big_number_size, bold=True,
                    color=pal.primary, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)
    enable_text_shrink(tb.text_frame)

    # Accent bar under the number
    bar_w = 2.0
    add_accent_bar(slide, (layout.slide_width_in - bar_w) / 2,
                   body_top + 0.3 + num_h + 0.05, bar_w, 0.06,
                   color=pal.accent)

    # Label
    tb = add_textbox(slide, layout.margin_left_in,
                     body_top + 0.3 + num_h + 0.2, width, 0.50)
    write_paragraph(tb.text_frame, stat_label,
                    size=typo.title_size, bold=True,
                    color=pal.primary_light, family=typo.family,
                    align=PP_ALIGN.CENTER, first=True)

    if context:
        tb = add_textbox(slide, 2.0, body_top + 0.3 + num_h + 0.9,
                         layout.slide_width_in - 4.0, 1.2)
        write_paragraph(tb.text_frame, context, size=typo.body_size + 2,
                        color=pal.text_dark, family=typo.family,
                        align=PP_ALIGN.CENTER, first=True)
    return slide


def add_chrome_title(slide, title, theme, page_number):
    add_title(slide, title, theme)
    add_footer(slide, theme, page_number=page_number)
